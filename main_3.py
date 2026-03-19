# -*- coding: utf-8 -*-
"""
Created on Fri Mar  6 10:01:36 2026

@tag: Xx_ScriptSniper_xX

@author: Albin
"""
import streamlit as st
# import pandas as pd
# import tempfile
import os
# import subprocess
# import sys
# import plotly.graph_objects as go
import json
import pickle
# from importlib import reload
from copy import deepcopy
# from libs.data_calculation import dataframe_modif
# from libs.Grapher import roundabout_specs, uturn_specs, stepsteer_specs, dlc_specs
from libs.figure_manager import FigureData, ManageData
from libs.color_manager import ColorManager
from libs.AA_utils import update_presets_json,load_presets,resend,respect,phunt
from libs.pptx_porter import pptX_tab
# from libs.spec_manager import SpecManager
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
# import os
from PIL import Image
from datetime import datetime
import shutil
import kaleido
kaleido.start_sync_server(n=5)


# @st.fragment
# def ppt_export_section():
    # st.subheader("Export PowerPoint Presentation")

    # # Maneuver selection checkboxes
    # chosen = {}
    # for m in st.session_state["tabs"].keys():
        # chosen[m] = st.checkbox(f"Maneuver {m}", value=True, key=f"chk_{m}")

    # # File settings
    # directory = st.text_input("Directory:", key="ppt_directory",
                              # value=r"C:\\project files\\General_PostP_tool\\scratch\\out")
    # filename = st.text_input("Name:", key="ppt_filename", value="Unyblyat")
    # template_path = st.text_input("Template path:", key="ppt_template",
                                  # value=r"C:\\project files\\General_PostP_tool\\scratch\\template.pptx")

    # col_main, col_settings = st.columns([3,1])
    # with col_settings:
        # with st.expander("⚙️ Settings", expanded=False):
            # cutfactor_1d = st.number_input("Cut factor 1D", min_value=0.1, max_value=5.0,
                                           # value=1.4, step=0.1, key="cutfactor_1d")
            # cutfactor_2d = st.number_input("Cut factor 2D", min_value=0.1, max_value=10.0,
                                           # value=7.5, step=0.1, key="cutfactor_2d")

    # # Export button
    # if st.button("Export", key="ppt_export"):
        # selected = [m for m, v in chosen.items() if v]
        # if not selected:
            # st.warning("Please select at least one maneuver.")
        # elif not directory or not filename or not template_path:
            # st.warning("Please provide directory, filename, and template path.")
        # else:
            # prs = Presentation(template_path)

            # # Create temporary folder for PNGs
            # temp_dir = os.path.join(directory, "ppt_temp_images")
            # os.makedirs(temp_dir, exist_ok=True)

            # # First slide (cover)
            # first_layout = prs.slide_layouts[0]
            # first_slide = prs.slides.add_slide(first_layout)
            # first_slide.placeholders[0].text = filename
            # today_str = datetime.today().strftime("%d %B %Y")
            # first_slide.placeholders[1].text = today_str

            # # Overview slide
            # overview_layout = prs.slide_layouts[2]
            # overview_slide = prs.slides.add_slide(overview_layout)
            # overview_slide.placeholders[11].text = filename
            # for i, maneuver_id in enumerate(selected[:8], start=1):
                # title_idx = 18 + (i - 1) * 2
                # num_idx = title_idx + 1
                # overview_slide.placeholders[title_idx].text = maneuver_id
                # overview_slide.placeholders[num_idx].text = f"{i:02d}"
            # for j in range(len(selected) + 1, 9):
                # title_idx = 18 + (j - 1) * 2
                # num_idx = title_idx + 1
                # for idx in (title_idx, num_idx):
                    # try:
                        # ph = overview_slide.placeholders[idx]
                        # ph.element.getparent().remove(ph.element)
                    # except IndexError:
                        # break

            # # Chapter title slides
            # for i, maneuver_id in enumerate(selected, start=1):
                # figures = st.session_state["tabs"][maneuver_id]["figures"]
                # title_layout = prs.slide_layouts[3]
                # title_slide = prs.slides.add_slide(title_layout)
                # title_slide.placeholders[14].text = f"{i:02d}"
                # title_slide.placeholders[15].text = maneuver_id

                # # Maneuver export slides
                # for idx, figdata in enumerate(figures):
                    # fig, split_ratio = figdata.render_page(
                        # page_data=figdata.page_data,
                        # vehicle_labels=st.session_state["impostors"].get(maneuver_id, {}).get(f"Page_{idx}", []),
                        # cutfactor_1d=cutfactor_1d,
                        # cutfactor_2d=cutfactor_2d
                    # )

                    # base_img_path = os.path.join(temp_dir, f"{maneuver_id}_Page{idx}_full.png")
                    # if isinstance(fig, plt.Figure):
                        # fig.savefig(base_img_path, dpi=300, bbox_inches="tight")
                        # plt.close(fig)
                    # else:
                        # fig.write_image(base_img_path, scale=2)

                    # img = Image.open(base_img_path)
                    # w, h = img.size
                    # plot_crop = img.crop((0, 0, int(w * split_ratio), h))
                    # legend_crop = img.crop((int(w * split_ratio), 0, w, h))
                    # plot_resized = respect(plot_crop, (854, 586))
                    # plot_path = os.path.join(temp_dir, f"{maneuver_id}_Page{idx}_plot.png")
                    # plot_resized.save(plot_path)
                    # legend_resized = resend(legend_crop, target_width_px=348, target_height_px=344)
                    # legend_path = os.path.join(temp_dir, f"{maneuver_id}_Page{idx}_legend.png")
                    # legend_resized.save(legend_path)

                    # slide_layout = prs.slide_layouts[4]
                    # slide = prs.slides.add_slide(slide_layout)
                    # try:
                        # slide.placeholders[11].text = f"{maneuver_id} - {figdata.page_data.get('title', f'Page{idx}')}"
                    # except KeyError:
                        # pass
                    # try:
                        # slide.placeholders[16].insert_picture(plot_path)
                    # except KeyError:
                        # slide.shapes.add_picture(plot_path, Inches(0.5), Inches(1.5),
                                                 # width=Inches(8.9), height=Inches(6.1))
                    # try:
                        # slide.placeholders[17].insert_picture(legend_path)
                    # except KeyError:
                        # slide.shapes.add_picture(legend_path, Inches(6.0), Inches(1.5))

            # Last_layout = prs.slide_layouts[9]
            # Last_slide = prs.slides.add_slide(Last_layout)

            # ppt_path = os.path.join(directory, f"{filename}.pptx")
            # prs.save(ppt_path)

            # # Clean up temporary folder
            # try:
                # shutil.rmtree(temp_dir)
            # except Exception as e:
                # st.warning(f"Could not delete temp folder: {e}")
            # st.success(f"PPT exported: {ppt_path}")



# from libs.AA_utils import process_data, render_page
# reload(SpecManager)
st.markdown('<div id="top"></div>', unsafe_allow_html=True)
# --- Load presets at startup ---
if "app_presets" not in st.session_state or "plot_presets" not in st.session_state:
    presets = load_presets()
    for k, v in presets.items():
        st.session_state[k] = v

st.set_page_config(page_title="Project PostProc", layout="wide")
st.title("Project PostProc")

# ======================
# Sidebar:
# ======================

st.sidebar.header("🎨 Theme & Plot Settings")
# --- Initialize applied theme if missing ---
if "applied_app_theme" not in st.session_state:
    st.session_state.applied_app_theme = st.session_state.app_presets["Preset 1"]

with st.sidebar.expander("🎨 App Theme Settings", expanded=False):
    selected_app_preset = st.radio(
        "Select App Preset",
        list(st.session_state.app_presets.keys()),
        key="selected_app_preset",
        horizontal=True
    )

    preset_values = st.session_state.app_presets[selected_app_preset]

    primary_color = st.color_picker("Primary color (buttons)", preset_values["primary"], key="tmp_app_primary")
    text_color = st.color_picker("Text color", preset_values["text"], key="tmp_app_text")
    background_color = st.color_picker("Background color", preset_values["background"], key="tmp_app_bg")
    secondary_background = st.color_picker("Sidebar background", preset_values["secondary"], key="tmp_app_secondary")

    col1, col2 = st.columns(2)
    with col1:
        if st.button("💾 Save App Theme"):
            st.session_state.app_presets[selected_app_preset] = {
                "primary": st.session_state.tmp_app_primary,
                "text": st.session_state.tmp_app_text,
                "background": st.session_state.tmp_app_bg,
                "secondary": st.session_state.tmp_app_secondary
            }
            update_presets_json({
                "app_presets": st.session_state.app_presets,
                "selected_app_preset": st.session_state.selected_app_preset,
            })
            st.toast("✅ App Theme Saved!", icon="🎉")

    with col2:
        if st.button("✅ Apply App Theme"):
            # Copy tmp values into preset
            st.session_state.app_presets[selected_app_preset] = {
                "primary": st.session_state.tmp_app_primary,
                "text": st.session_state.tmp_app_text,
                "background": st.session_state.tmp_app_bg,
                "secondary": st.session_state.tmp_app_secondary
            }
            # Update applied theme snapshot
            st.session_state.applied_app_theme = st.session_state.app_presets[selected_app_preset]
            update_presets_json({
                "app_presets": st.session_state.app_presets,
                "selected_app_preset": st.session_state.selected_app_preset,
                "applied_app_theme": st.session_state.app_presets[selected_app_preset]
            })
            st.toast("✅ App Theme Applied!", icon="🎉")

# --- Always inject CSS from applied theme snapshot ---
ui_colors = st.session_state.applied_app_theme
global_color_manager = ColorManager(ui_colors=ui_colors)

st.markdown(
    f"""
    <style>
    /* ---------------- GLOBAL ---------------- */
    .stApp {{
        background-color: {global_color_manager.background};
        color: {global_color_manager.text};
    }}

    /* Headers */
    h1, h2, h3, h4, h5, h6 {{
        color: {global_color_manager.text} !important;
    }}

    /* Widget labels & text */
    label, .stMarkdown, .stText, .stRadio, .stCheckbox, .stSelectbox {{
        color: {global_color_manager.text} !important;
    }}
    /* Checkbox label text rendered via Markdown container */
    div[data-testid="stMarkdownContainer"] p {{
        color: {global_color_manager.text} !important;
    }}
    /* ---------------- SIDEBAR ---------------- */
    .stSidebar {{
        background-color: {global_color_manager.secondary};
        color: {global_color_manager.text};
    }}
    .stSidebar, .stSidebar div, .stSidebar p, .stSidebar label {{
        color: {global_color_manager.text} !important;
    }}

    /* Sidebar collapse/expand toggle */
    span.st-emotion-cache-189uypx {{
        background-color: {global_color_manager.secondary} !important;
        border-radius: 4px;
        padding: 4px;
        display: inline-flex;
        align-items: center;
        justify-content: center;
    }}
    span.st-emotion-cache-189uypx span[data-testid="stIconMaterial"] {{
        color: {global_color_manager.text} !important;
    }}
    span.st-emotion-cache-189uypx:hover {{
        background-color: {global_color_manager.primary} !important;
    }}

    /* ---------------- HEADER ---------------- */
    header[data-testid="stHeader"] {{
        background-color: {global_color_manager.background};
        color: {global_color_manager.text};
    }}
    header[data-testid="stHeader"] * {{
        color: {global_color_manager.text} !important;
    }}

    /* ---------------- BUTTONS ---------------- */
    div.stButton > button:first-child {{
        background-color: {global_color_manager.primary};
        color: {global_color_manager.text};
    }}

    /* ---------------- SELECTBOX ---------------- */
    div[data-baseweb="select"] > div {{
        background-color: {global_color_manager.secondary};
        color: {global_color_manager.text};
        border-radius: 4px;
    }}
    .stSelectbox div[data-baseweb="select"] span {{
        color: {global_color_manager.text} !important;
    }}

    /* ---------------- FILE UPLOADER ---------------- */
    div[data-testid="stFileUploader"] {{
        background-color: {global_color_manager.secondary};
        border-radius: 6px;
        padding: 8px;
    }}
    div[data-testid="stFileUploader"] button {{
        background-color: {global_color_manager.primary};
        color: {global_color_manager.text};
        border: none;
        border-radius: 4px;
    }}
    div[data-testid="stFileUploader"] section[data-testid="stFileUploaderDropzone"] button[data-testid="stBaseButton-secondary"]:hover {{
        background-color: {global_color_manager.primary};
        color: {global_color_manager.text};
    }}

    /* Dropzone background */
    section[data-testid="stFileUploaderDropzone"] {{
        background-color: {global_color_manager.background} !important;
        border: 2px dashed {global_color_manager.primary} !important;
        border-radius: 6px;
        color: {global_color_manager.text} !important;
    }}
    /* Dropzone instructions text */
    div[data-testid="stFileUploaderDropzoneInstructions"] span {{
        color: {global_color_manager.text} !important;
    }}
    div[data-testid="stFileUploaderDropzoneInstructions"] span:last-child {{
        color: {global_color_manager.text} !important;
    }}
    div[data-testid="stFileUploaderDropzoneInstructions"] svg {{
        fill: {global_color_manager.text} !important;
    }}
    /* ---------------- TOAST ---------------- */
    div.stToast {{
        background-color: {global_color_manager.secondary} !important;
        color: {global_color_manager.text} !important;
        border-radius: 6px;
        padding: 8px;
    }}
    div.stToast p {{
        color: {global_color_manager.text} !important;
        font-weight: bold;
    }}
    /* ---------------- EXPANDERS ---------------- */
    details > summary {{
        background-color: {global_color_manager.secondary} !important;
        color: {global_color_manager.text} !important;
        border-radius: 4px;
        padding: 4px 8px;
    }}
    details > summary:hover {{
        background-color: {global_color_manager.primary} !important;
    }}
    </style>
    """,
    unsafe_allow_html=True
)

if "applied_plot_theme" not in st.session_state:
    st.session_state.applied_plot_theme = st.session_state.plot_presets["Preset 1"]

# --- Expander 2: Plot Theme ---
with st.sidebar.expander("📊 Plot Theme Settings", expanded=False):
    selected_plot_preset = st.radio(
        "Select Plot Preset",
        list(st.session_state.plot_presets.keys()),
        index=list(st.session_state.plot_presets.keys()).index(
            st.session_state.get("selected_plot_preset", "Preset 1")
        ),
        horizontal=True
    )

    preset_values = st.session_state.plot_presets[selected_plot_preset]

    # Color pickers (UI only)
    plot_title_color = st.color_picker("Plot title color", preset_values.get("plot_title", "#FFFFFF"), key="plot_title")
    axis_title_color = st.color_picker("Axis title color", preset_values["axis_title"], key="plot_axis_title")
    legend_text_color = st.color_picker("Legend text color", preset_values["legend_text"], key="plot_legend_text")
    tick_label_color = st.color_picker("Tick label color", preset_values["tick_label"], key="plot_tick_label")
    plot_bgcolor = st.color_picker("Plot background", preset_values["plot_bg"], key="plot_bg")
    grid_color = st.color_picker("Gridline color", preset_values["grid"], key="plot_grid")

    # --- 16-color palette (two rows of 8) ---
    st.markdown("**Plot Line Color Palette**")
    palette_colors = []
    for row in range(2):
        cols = st.columns(8)
        for i in range(8):
            idx = row * 8 + i
            with cols[i]:
                color = st.color_picker(
                    f"L{idx+1}",
                    preset_values.get(f"plot_line_color{idx+1}", "#000000"),
                    key=f"plot_line_color{idx+1}"
                )
                palette_colors.append(color)

    use_palette_key = f"{selected_plot_preset}_use_palette"
    if use_palette_key not in st.session_state:
        st.session_state[use_palette_key] = preset_values.get("use_palette", False)
    use_palette = st.checkbox("Use palette cycling", value=st.session_state[use_palette_key], key=use_palette_key)

    col1, col2 = st.columns(2)
    with col1:
        if st.button("💾 Save Plot Theme"):
            st.session_state.plot_presets[selected_plot_preset] = {
                "axis_title": axis_title_color,
                "tick_label": tick_label_color,
                "legend_text": legend_text_color,
                "plot_title": plot_title_color,
                "plot_bg": plot_bgcolor,
                "grid": grid_color,
                **{f"plot_line_color{i+1}": palette_colors[i] for i in range(16)},
                "use_palette": st.session_state[use_palette_key]
            }
            st.session_state.selected_plot_preset = selected_plot_preset
            update_presets_json({
                "plot_presets": st.session_state.plot_presets,
                "selected_plot_preset": selected_plot_preset
            })
            st.toast("✅ Saved Plot Themes!", icon="🎉")
    
    with col2:
        if st.button("✅ Apply Plot Theme"):
            st.session_state.applied_plot_theme = {
                "axis_title": axis_title_color,
                "tick_label": tick_label_color,
                "legend_text": legend_text_color,
                "plot_title": plot_title_color,
                "plot_bg": plot_bgcolor,
                "grid": grid_color,
                **{f"plot_line_color{i+1}": palette_colors[i] for i in range(16)},
                "use_palette": st.session_state[use_palette_key]
            }
            st.session_state.selected_plot_preset = selected_plot_preset
            update_presets_json({
                "plot_presets": st.session_state.plot_presets,
                "selected_plot_preset": selected_plot_preset,
                "applied_plot_theme": st.session_state.applied_plot_theme
            })
            st.toast("✅ Applied Plot Themes!", icon="🎉")

# After Apply Plot Theme
plot_colors = st.session_state.applied_plot_theme
# Initialize global ColorManager
global_color_manager = ColorManager(ui_colors, plot_colors)

st.sidebar.header("📂 Inputs")


# Store uploaded file handles in session_state
if "datasets" not in st.session_state:
    st.session_state["datasets"] = {}

# --- File upload ---
uploaded_files = st.sidebar.file_uploader("Upload simulation files", type=["json"], accept_multiple_files=True)

if uploaded_files:
    for file in uploaded_files:
        file_name = file.name
        data = json.load(file)
        data["__file_name__"] = file_name
        st.session_state["datasets"][file_name] = data


# --- Alias editing UI ---
if st.session_state.get("datasets"):
    with st.sidebar.expander("Aliases", expanded=False):
        for idx, file_name in enumerate(st.session_state["datasets"].keys()):
            alias_key = f"alias_{idx}_{file_name}"
            # default_val = st.session_state.get("aliases", {}).get(file_name, file_name)
            # if alias_key not in st.session_state:
                # st.session_state[alias_key] = default_val

            # Store alias directly in aliases dict
            new_alias = st.text_input(f"{file_name}", key=alias_key)
            st.session_state.setdefault("aliases", {})
            st.session_state["aliases"][file_name] = new_alias

        if st.button("Update Aliases"):
            st.toast("✅ Aliases updated successfully!", icon="🎉")


# --- Initialize managers once ---
if "managedata" not in st.session_state:
    st.session_state["managedata"] = ManageData()

# --- Button trigger ---
if st.sidebar.button("🚀 Let's Roll", use_container_width=True):
    if not uploaded_files:
        st.toast("⚠️ Please upload at least one JSON file.", icon="❗")
    else:
        try:
            data_inputs = [raw_file for raw_file in uploaded_files]
            format_map = {
                "DA80DAGene": "DA_Gene_v3.json",
                "DA80DAGeney": "DA_Gene_v4.json",
                "YASRMultiplicator" : "YASRMultiplicator.json",
                "Throttle_off_STD_v2a" : "Cornering_Throttle_off_STD.json",
                "DA80_Kph_LH" : "DA_80kph_new.json",
            }
            BASE_DIR = os.path.dirname(os.path.abspath(__file__))
            format_folder = os.path.join(BASE_DIR, "ReportManager")

            # Build alias mapping from session state
            aliases_map = st.session_state.get("aliases", {})
            # Step 1: classify, now passing aliases
            data_inputs = list(st.session_state["datasets"].values())
            classified_output, impostors, stowaways = st.session_state["managedata"].process_data(
                data_inputs,
                format_map,
                format_folder=format_folder,
                aliases=aliases_map
            )
            phunt(impostors,name="imps")

            # Save both outputs into session state
            st.session_state["classified_output"] = classified_output
            st.session_state["impostors"] = impostors   # <-- this was missing
            st.session_state["stowaways"] = stowaways 
            with open("classified_output.pkl", "wb") as f:
                pickle.dump((classified_output,impostors), f)

            new_figures = []

            # Step 2: build FigureData per page
            new_tabs = {}

            # Step 2: build FigureData per page
            for maneuver_id, pages in classified_output.items():
                if "error" in pages:
                    continue
                maneuver_figs = []
                for page_name, page_data in pages.items():
                    if not isinstance(page_data, dict) or "subgraphs" not in page_data:
                        continue
                    figdata = FigureData(page_data, color_manager=deepcopy(global_color_manager))
                    page_data = figdata.build_spec()
                    figdata.page_data = page_data
                    maneuver_figs.append(figdata)

                new_tabs[maneuver_id] = {"figures": maneuver_figs}

            st.session_state["tabs"] = new_tabs
        except Exception as e:
            st.error("Invalid data")
            st.exception(e)

# --- Outer loop: apply SpecManager + render ---
if "tabs" in st.session_state and "impostors" in st.session_state:
    tabs = st.tabs(list(st.session_state["tabs"].keys()) + ["pptX"])

    # --- Maneuver tabs ---
    for tab, maneuver_id in zip(tabs[:-1], st.session_state["tabs"].keys()):
        with tab:
            st.subheader(f"Simulation Results for: {maneuver_id}")
            for idx, figdata in enumerate(st.session_state["tabs"][maneuver_id]["figures"], start=0):
                # try:
                #     spec_manager = SpecManager(figdata)
                #     spec_manager.build_ui()
                #     spec_manager.apply_updates()
                #     figdata = spec_manager.figdata
                # except TypeError:
                #     st.warning("No Plot theme applied. Please select a color(s) manually.")

                page_name = figdata.page_data.get("title", f"Page{idx}")
                vehicle_labels = st.session_state["impostors"].get(maneuver_id, {}).get(f"Page_{idx}", [])
                fig,dump = figdata.render_page(page_data=figdata.page_data, vehicle_labels=vehicle_labels)

                if isinstance(fig, plt.Figure):
                    st.pyplot(fig)
                else:
                    # 🔹 Add a unique key so relayoutData is tracked
                    st.plotly_chart(fig, use_container_width=True, key=f"chart_{maneuver_id}_{idx}")

                    # Capture zoom/pan state
                    # relayout = st.session_state.get(f"chart_{maneuver_id}_{idx}", {}).get("relayoutData", {})
                    # print(relayout)
                    # if "xaxis.range[0]" in relayout and "xaxis.range[1]" in relayout:
                        # st.session_state[f"zoom_x_{maneuver_id}_{idx}"] = (
                            # relayout["xaxis.range[0]"], relayout["xaxis.range[1]"]
                        # )
                    # if "yaxis.range[0]" in relayout and "yaxis.range[1]" in relayout:
                        # st.session_state[f"zoom_y_{maneuver_id}_{idx}"] = (
                            # relayout["yaxis.range[0]"], relayout["yaxis.range[1]"]
                        # )

                st.markdown("<br><br>", unsafe_allow_html=True)
    if "default_selected" not in st.session_state:
        default_mapping = {
            "DA80_Kph_LH": {
                "0": ["DA_80kph_LH__linear_slope_understeer"],
                "3": ["DA_80kph_LH__linear_slope_spe_rear_slip"],
                "4": ["DA_80kph_LH__characteristic_speed"],
            },
            "DA80DAGene": {
                "5": ["DA80DAGene_DAGeneWob__accy_stw_response_1Hz"],
                "17": ["DA80DAGene_DAGeneWob__yaw_stw_response_1Hz"],
                "19": ["DA80DAGene_DAGeneWob__roll_stw_response_1Hz"],
                "20": ["DA80DAGene_DAGeneWob__accy_yaw_response_1Hz"],
            },
        }
    
        default_selected = []
    
        # Only add defaults if the maneuver/page/stowaway exists in current session
        for maneuver, pages in default_mapping.items():
            if maneuver in st.session_state["tabs"]:  # check maneuver exists
                for page_no, stowaways in pages.items():
                    # check page index is valid
                    if int(page_no) < len(st.session_state["tabs"][maneuver]["figures"]):
                        # check stowaway keys exist
                        stowaway_dict = st.session_state.get("stowaways", {}).get(maneuver, {})
                        if isinstance(stowaway_dict, dict):
                            for st_key in stowaways:
                                if st_key in stowaway_dict:
                                    val = f"man{maneuver}_page{page_no}_{st_key}"
                                    default_selected.append(val)
                        elif isinstance(stowaway_dict, list):
                            for st_key in stowaways:
                                if st_key in stowaway_dict:
                                    val = f"man{maneuver}_page{page_no}_{st_key}"
                                    default_selected.append(val)
    
        st.session_state["default_selected"] = default_selected
    # --- PPT export tab ---
    # Call inside your tab loop
    with tabs[-1]:
        pptX_tab()
        # st.button("update", key="onckle")

# Place this once in your script (anywhere after st imports)
st.markdown("""
    <style>
    html {
        scroll-behavior: smooth; /* enables smooth scroll */
    }
    #goTopBtn {
        position: fixed;
        bottom: 40px;
        right: 40px;
        z-index: 999;
        border: none;
        outline: none;
        background-color: #4CAF50;
        color: white;
        cursor: pointer;
        padding: 12px 18px;
        border-radius: 6px;
        font-size: 14px;
        box-shadow: 0 2px 6px rgba(0,0,0,0.3);
    }
    #goTopBtn:hover {
        background-color: #45a049;
    }
    </style>

    <a href="#top"><button id="goTopBtn">⬆ Top</button></a>
""", unsafe_allow_html=True)



