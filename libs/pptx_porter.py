# -*- coding: utf-8 -*-
"""
Created on Tue Mar 17 10:13:51 2026

@tag: Xx_ScriptSniper_xX

@author: Albin
"""
from pptx import Presentation
from pptx.util import Inches
import shutil
from PIL import Image
from datetime import datetime
import streamlit as st
import os
import matplotlib.pyplot as plt
from libs.AA_utils import resend,respect,phunt
from st_ant_tree import st_ant_tree
from pptx.util import Cm
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
import re
from pptx.dml.color import RGBColor
from pptx.util import Pt

def debug_placeholders(slide):
    """
    Print out all placeholder indexes and names for a given slide.
    Helps you identify which placeholder to use for tables.
    """
    for ph in slide.placeholders:
        print(f"Index: {ph.placeholder_format.idx}, Type: {ph.name}")

def maneuver_selection(prefix=""):
    chosen = {}
    for m in st.session_state["tabs"].keys():
        chosen[m] = st.checkbox(f"Maneuver {m}", value=True, key=f"{prefix}chk_{m}")
    return chosen

def file_settings(prefix=""):
    directory = st.text_input("Directory:", key=f"{prefix}ppt_directory",
                              value=r"C:\\project files\\General_PostP_tool\\scratch\\out")
    filename = st.text_input("Name:", key=f"{prefix}ppt_filename", value="Unyblyat")
    # template_path = st.text_input("Template path:", key=f"{prefix}ppt_template",
    #                               value=r"C:\\project files\\General_PostP_tool\\scratch\\template.pptx")
    base_dir = os.path.dirname(__file__)
    template_path = os.path.join(base_dir, "template.pptx")
    return directory, filename, template_path

def export_settings(prefix=""):
    with st.expander("⚙️ Settings", expanded=False):
        cutfactor_1d = st.number_input("Cut factor 1D", min_value=0.1, max_value=5.0,
                                       value=1.4, step=0.1, key=f"{prefix}cutfactor_1d")
        cutfactor_2d = st.number_input("Cut factor 2D", min_value=0.1, max_value=10.0,
                                       value=7.5, step=0.1, key=f"{prefix}cutfactor_2d")
    return cutfactor_1d, cutfactor_2d

# --- Export helpers ---
def create_cover_slide(prs, filename):
    first_layout = prs.slide_layouts[0]
    first_slide = prs.slides.add_slide(first_layout)
    first_slide.placeholders[0].text = filename
    today_str = datetime.today().strftime("%d %B %Y")
    first_slide.placeholders[1].text = today_str

def create_overview_slide(prs, selected, filename):
    overview_layout = prs.slide_layouts[2]
    overview_slide = prs.slides.add_slide(overview_layout)
    overview_slide.placeholders[11].text = filename
    for i, maneuver_id in enumerate(selected[:8], start=1):
        title_idx = 18 + (i - 1) * 2
        num_idx = title_idx + 1
        overview_slide.placeholders[title_idx].text = maneuver_id
        overview_slide.placeholders[num_idx].text = f"{i:02d}"
    for j in range(len(selected) + 1, 9):
        title_idx = 18 + (j - 1) * 2
        num_idx = title_idx + 1
        for idx in (title_idx, num_idx):
            try:
                ph = overview_slide.placeholders[idx]
                ph.element.getparent().remove(ph.element)
            except IndexError:
                break

def longest_common_prefix(strings):
    if not strings: return ""
    s1, s2 = min(strings), max(strings)
    for i, c in enumerate(s1):
        if c != s2[i]:
            return s1[:i]
    return s1

def longest_common_suffix(strings):
    if not strings: return ""
    rev = [s[::-1] for s in strings]
    return longest_common_prefix(rev)[::-1]

def clean_vehicle_names(vehicle_names):
    """Remove common prefix and suffix across all vehicle names and truncate to 34 chars."""
    cleaned = {}
    if len(vehicle_names) > 1:
        prefix = longest_common_prefix(vehicle_names)
        suffix = longest_common_suffix(vehicle_names)
    else:
        prefix = ""
        suffix = ""

    for name in vehicle_names:
        new_name = name
        if prefix:
            new_name = new_name[len(prefix):]
        if suffix and new_name.endswith(suffix):
            new_name = new_name[:-len(suffix)]

        new_name = new_name.strip("_- ")

        if len(new_name) > 33:
            new_name = new_name[:31] + "..."

        cleaned[name] = new_name
    return cleaned


def shorten_st_key(st_key: str, maneuver_id: str) -> str:
    """Remove maneuver_id and trailing underscores/dashes, case-insensitively."""
    pattern = '^' + re.escape(maneuver_id) + r'[_\-]*'
    # Adding re.IGNORECASE here fixes the mismatch
    return re.sub(pattern, '', st_key, flags=re.IGNORECASE)

def insert_stowaway_tables(slide, maneuver_id, page_idx, table_mapping):
    page_key = str(page_idx)
    if maneuver_id in table_mapping and page_key in table_mapping[maneuver_id]:
        for st_key in table_mapping[maneuver_id][page_key]:
            stowaway_dict = st.session_state.get("stowaways", {}).get(maneuver_id, {})
            if st_key in stowaway_dict:
                vehicle_data = stowaway_dict[st_key]

                rows = len(vehicle_data) + 1
                cols = 2

                ph = slide.placeholders[18]  # adjust index to your layout
                graphic_frame = ph.insert_table(rows, cols)
                table = graphic_frame.table

                # --- Override style to "No Style, No Grid" ---
                no_style_no_grid = '{21E7435A-961F-4814-916E-73D13E9FBCD9}'
                tbl = graphic_frame._element.graphic.graphicData.tbl
                tbl[0][-1].text = no_style_no_grid

                # --- Header row ---
                table.cell(0, 0).text = "Vehicle"
                table.cell(0, 1).text = shorten_st_key(st_key, maneuver_id)
                for c in range(cols):
                    cell = table.cell(0, c)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(220, 220, 220)
                    para = cell.text_frame.paragraphs[0]
                    run = para.runs[0]
                    run.font.size = Pt(10)
                    run.font.bold = True
                    para.alignment = PP_ALIGN.CENTER
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    cell.margin_left = Pt(2)
                    cell.margin_top = Pt(0)
                    cell.margin_bottom = Pt(0)
                    cell.margin_right = Pt(0)

                # --- Clean vehicle names ---
                cleaned_map = clean_vehicle_names(list(vehicle_data.keys()))

                # --- Data rows ---
                for r, (veh, val) in enumerate(vehicle_data.items(), start=1):
                    table.cell(r, 0).text = cleaned_map[veh]
                    try:
                        num = float(val)
                        if abs(num) < 0.01:
                            table.cell(r, 1).text = f"{num:.4f}"
                        else:
                            table.cell(r, 1).text = f"{num:.2f}"
                    except (ValueError, TypeError):
                        table.cell(r, 1).text = str(val)

                    for c in range(cols):
                        cell = table.cell(r, c)
                        para = cell.text_frame.paragraphs[0]
                        if para.runs:
                            run = para.runs[0]
                        else:
                            run = para.add_run()  # create a run if none exist
                        run.font.size = Pt(10)
                        cell.text_frame.word_wrap = False
                        cell.margin_left = Pt(2)
                        cell.margin_top = Pt(0)
                        cell.margin_bottom = Pt(0)
                        cell.margin_right = Pt(0)

                        if c == 1:
                            para.alignment = PP_ALIGN.CENTER

                # --- Column widths ---
                total_width = graphic_frame.width
                table.columns[0].width = int(total_width * 0.7)
                table.columns[1].width = int(total_width * 0.3)

                # --- Compress row heights ---
                total_height = graphic_frame.height
                natural_height = total_height / rows

                # Define max row height for 8 vehicles ≈ 0.45 cm
                max_row_height = Cm(0.45)

                # Scale: if fewer rows, allow taller rows; if more, cap at max_row_height
                row_height = min(natural_height, max_row_height)

                for row in table.rows:
                    row.height = row_height

def create_maneuver_slides(prs, maneuver_id, figures, cutfactor_1d, cutfactor_2d, temp_dir,
                           table_mapping=None, slide_number=None):
    title_layout = prs.slide_layouts[3]
    title_slide = prs.slides.add_slide(title_layout)
    if slide_number is not None:
        title_slide.placeholders[14].text = f"{slide_number:02d}"
    title_slide.placeholders[15].text = maneuver_id

    for idx, figdata in enumerate(figures):
        fig, split_ratio = figdata.render_page(
            page_data=figdata.page_data,
            vehicle_labels=st.session_state["impostors"].get(maneuver_id, {}).get(f"Page_{idx}", []),
            cutfactor_1d=cutfactor_1d,
            cutfactor_2d=cutfactor_2d
        )

        base_img_path = os.path.join(temp_dir, f"{maneuver_id}_Page{idx}_full.png")
        if isinstance(fig, plt.Figure):
            fig.savefig(base_img_path, dpi=300, bbox_inches="tight")
            plt.close(fig)
        else:
            fig.write_image(base_img_path, scale=2)
            img = Image.open(base_img_path)
            w, h = img.size
            margin = 50
            figcrop = img.crop((margin, 0, w, h))
            figcrop.save(base_img_path)

        img = Image.open(base_img_path)
        w, h = img.size
        plot_crop = img.crop((0, 0, int(w * split_ratio), h))
        legend_crop = img.crop((int(w * split_ratio), 0, w, h))
        plot_resized = respect(plot_crop, (854, 586))
        plot_path = os.path.join(temp_dir, f"{maneuver_id}_Page{idx}_plot.png")
        plot_resized.save(plot_path)
        legend_resized = resend(legend_crop, target_width_px=348, target_height_px=344)
        legend_path = os.path.join(temp_dir, f"{maneuver_id}_Page{idx}_legend.png")
        legend_resized.save(legend_path)
        
        slide_layout = prs.slide_layouts[4]
        slide = prs.slides.add_slide(slide_layout)
        try:
            slide.placeholders[11].text = f"{maneuver_id} - {figdata.page_data.get('title', f'Page{idx}')}"
        except KeyError:
            pass
        try:
            slide.placeholders[16].insert_picture(plot_path)
        except KeyError:
            slide.shapes.add_picture(plot_path, Inches(0.5), Inches(1.5),
                                     width=Inches(8.9), height=Inches(6.1))
        try:
            slide.placeholders[17].insert_picture(legend_path)
        except KeyError:
            slide.shapes.add_picture(legend_path, Inches(6.0), Inches(1.5))

        # --- NEW: insert stowaway tables if mapped ---
        # if table_mapping:
        #     insert_stowaway_tables(slide, maneuver_id, idx, table_mapping)
            
        if maneuver_id in table_mapping and str(idx) in table_mapping[maneuver_id]:
            insert_stowaway_tables(slide, maneuver_id, idx, table_mapping)
        else:
            # remove unused table placeholder if no data
            try:
                ph = slide.placeholders[18]  # adjust index to your layout
                ph.element.getparent().remove(ph.element)
            except IndexError:
                pass
# @st.fragment
# def pptX_tree(key=None):
#     st.subheader("Export PowerPoint Presentation (Tree)")

#     tree_data = [
#         {"title": "cover_slide", "value": "cover"},
#         {"title": "overview_slide", "value": "overview"},
#     ]

#     for m in st.session_state["tabs"].keys():
#         man_node = {"title": f"Maneuver {m}", "value": f"man{m}", "children": []}
#         for page_idx, figdata in enumerate(st.session_state["tabs"][m]["figures"], start=0):
#             page_name = figdata.page_data.get("title", f"Page{page_idx}")
#             page_node = {"title": page_name, "value": f"man{m}_page{page_idx}", "children": []}

#             stowaway_dict = st.session_state.get("stowaways", {}).get(m, None)
#             if stowaway_dict is None:
#                 stowaway_dict = st.session_state["tabs"][m].get("stowaway", {})

#             if isinstance(stowaway_dict, dict):
#                 for option_key in stowaway_dict.keys():
#                     page_node["children"].append({
#                         "title": option_key,
#                         "value": f"man{m}_page{page_idx}_{option_key}"
#                     })
#             elif isinstance(stowaway_dict, list):
#                 for param in stowaway_dict:
#                     page_node["children"].append({
#                         "title": param,
#                         "value": f"man{m}_page{page_idx}_{param}"
#                     })

#             man_node["children"].append(page_node)

#         tree_data.append(man_node)

#     selected_values = st_ant_tree(
#         treeData=tree_data,
#         showSearch=True,
#         treeCheckable=True,
#         defaultValue=st.session_state["default_selected"],
#         # key=key,
#         max_height=600,
#         placeholder="Search and select"
#     )

#     table_mapping = {}

#     if selected_values:
#         for val in selected_values:
#             if "_page" in val:
#                 man_part, rest = val.split("_page", 1)
#                 page_part, stowaway_key = rest.split("_", 1)

#                 maneuver_id = man_part.replace("man", "", 1)
#                 page_no = page_part

#                 if maneuver_id not in table_mapping:
#                     table_mapping[maneuver_id] = {}
#                 if page_no not in table_mapping[maneuver_id]:
#                     table_mapping[maneuver_id][page_no] = []

#                 table_mapping[maneuver_id][page_no].append(stowaway_key)

#     st.session_state["table_mapping"] = table_mapping
#     st.write("Table mapping:", table_mapping)

@st.fragment
def pptX_tree(key=None):
    st.subheader("Export PowerPoint Presentation (Tree)")

    if "show_tree" not in st.session_state:
        st.session_state["show_tree"] = False

    if st.button("➕ Add Table", key="add_table_btn"):
        st.session_state["show_tree"] = True

    if st.session_state["show_tree"]:
        tree_data = [
            {"title": "cover_slide", "value": "cover"},
            {"title": "overview_slide", "value": "overview"},
        ]

        for m in st.session_state["tabs"].keys():
            man_node = {"title": f"Maneuver {m}", "value": f"man{m}", "children": []}
            for page_idx, figdata in enumerate(st.session_state["tabs"][m]["figures"], start=0):
                page_name = figdata.page_data.get("title", f"Page{page_idx}")
                page_node = {"title": page_name, "value": f"man{m}_page{page_idx}", "children": []}

                stowaway_dict = st.session_state.get("stowaways", {}).get(m, None)
                if stowaway_dict is None:
                    stowaway_dict = st.session_state["tabs"][m].get("stowaway", {})

                if isinstance(stowaway_dict, dict):
                    for option_key in stowaway_dict.keys():
                        page_node["children"].append({
                            "title": option_key,
                            "value": f"man{m}_page{page_idx}_{option_key}"
                        })
                elif isinstance(stowaway_dict, list):
                    for param in stowaway_dict:
                        page_node["children"].append({
                            "title": param,
                            "value": f"man{m}_page{page_idx}_{param}"
                        })

                man_node["children"].append(page_node)

            tree_data.append(man_node)
        
        selected_values = st_ant_tree(
            treeData=tree_data,
            showSearch=True,
            treeCheckable=True,
            defaultValue=st.session_state["default_selected"],
            key=key,
            max_height=600,
            placeholder="Search and select"
        )

        table_mapping = {}

        if selected_values:
            for val in selected_values:
                if "_page" in val:
                    man_part, rest = val.split("_page", 1)
                    page_part, stowaway_key = rest.split("_", 1)

                    maneuver_id = man_part.replace("man", "", 1)
                    page_no = page_part

                    if maneuver_id not in table_mapping:
                        table_mapping[maneuver_id] = {}
                    if page_no not in table_mapping[maneuver_id]:
                        table_mapping[maneuver_id][page_no] = []

                    table_mapping[maneuver_id][page_no].append(stowaway_key)
        
        st.session_state["table_mapping"] = table_mapping
        st.write("Table mapping:", table_mapping)

@st.fragment
def pptX_export():
    st.subheader("Export PowerPoint Presentation (Classic)")
    chosen = maneuver_selection(prefix="classic_")
    directory, filename, template_path = file_settings(prefix="classic_")
    cutfactor_1d, cutfactor_2d = export_settings(prefix="classic_")

    if st.button("Export", key="ppt_export"):
        selected = [m for m, v in chosen.items() if v]
        if not selected:
            st.warning("Please select at least one maneuver.")
        elif not directory or not filename or not template_path:
            st.warning("Please provide directory, filename, and template path.")
        else:
            prs = Presentation(template_path)
            temp_dir = os.path.join(directory, "ppt_temp_images")
            os.makedirs(temp_dir, exist_ok=True)

            create_cover_slide(prs, filename)
            create_overview_slide(prs, selected, filename)
            
            for i, maneuver_id in enumerate(selected, start=1):
                figures = st.session_state["tabs"][maneuver_id]["figures"]
                table_mapping = st.session_state.get("table_mapping", {})
                create_maneuver_slides(prs, maneuver_id, figures,
                                       cutfactor_1d, cutfactor_2d, temp_dir,
                                       slide_number=i,
                                       table_mapping = table_mapping)

            prs.slides.add_slide(prs.slide_layouts[9])
            ppt_path = os.path.join(directory, f"{filename}.pptx")
            prs.save(ppt_path)
            try:
                shutil.rmtree(temp_dir)
            except Exception as e:
                st.warning(f"Could not delete temp folder: {e}")
            st.success(f"PPT exported: {ppt_path}")

# --- Tabs in the UI ---
def pptX_tab():
    tab1, tab2 = st.tabs(["MapDeTree", "Export"])
    with tab1:
        pptX_tree(key="pptx_tree")
    with tab2:
        pptX_export()   # unique key here


