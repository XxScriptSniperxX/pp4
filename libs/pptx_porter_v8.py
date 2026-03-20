# -*- coding: utf-8 -*-
"""
Created on Tue Mar 17 10:13:51 2026

@tag: Xx_ScriptSniper_xX
@author: Albin

PPTX PORTER V5 - Production Ready
- HIGH QUALITY: Matplotlib DPI=300, Plotly scale=2
- cv2 LEGEND DETECTION: Smart contour detection with padding
- FLEXIBLE RESIZING: Separate sizes for 1D/2D plots
- STREAMLIT CLOUD: No file dependencies, BytesIO throughout
"""
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import cv2
import numpy as np
import matplotlib.pyplot as plt
import streamlit as st
import os
import re
import tempfile
import io
from datetime import datetime
from st_ant_tree import st_ant_tree

# ============================================================================
# CONFIGURATION & CONSTANTS
# ============================================================================

@st.cache_resource
def get_temp_dir():
    """Get temp directory for Streamlit Cloud."""
    temp_base = tempfile.gettempdir()
    temp_dir = os.path.join(temp_base, "streamlit_pptx_export")
    os.makedirs(temp_dir, exist_ok=True)
    return temp_dir


@st.cache_resource
def get_template_path() -> str:
    """Get template path from package directory."""
    base_dir = os.path.dirname(os.path.abspath(__file__))
    template = os.path.join(base_dir, "template.pptx")
    if not os.path.exists(template):
        raise FileNotFoundError(f"Template not found: {template}")
    return template


# ============================================================================
# CONSTANTS - QUALITY & SIZING
# ============================================================================

# Plot sizing
PLOT_SIZE_1D = (854, 586)           # 1D Matplotlib plots
PLOT_SIZE_2D = (854, 586)           # 2D Plotly plots (larger to avoid squeezing)
LEGEND_SIZE = (348, 344)

# Legend cropping padding (adjust to show/hide borders)
LEGEND_PADDING = {
    'left': 15,                      # Show left border
    'right': 15,                     # Show right border
    'top': 15,                       # Show top border
    'bottom': 15                     # Show bottom border
}

# Quality settings
MATPLOTLIB_DPI = 300                # HIGH QUALITY
PLOTLY_SCALE = 2                    # HIGH QUALITY
LEGEND_FALLBACK_RATIO = 0.75

# Tables
DEFAULT_MAX_LEGEND_HEIGHT = Cm(0.45)
TABLE_HEADER_COLOR = RGBColor(220, 220, 220)
TABLE_FONT_SIZE = Pt(10)

# ============================================================================
# UTILITY FUNCTIONS
# ============================================================================

def longest_common_prefix(strings: list[str]) -> str:
    """Find longest common prefix."""
    if not strings:
        return ""
    s1, s2 = min(strings), max(strings)
    for i, c in enumerate(s1):
        if c != s2[i]:
            return s1[:i]
    return s1


def longest_common_suffix(strings: list[str]) -> str:
    """Find longest common suffix."""
    if not strings:
        return ""
    rev = [s[::-1] for s in strings]
    return longest_common_prefix(rev)[::-1]


def shorten_st_key(st_key: str, maneuver_id: str) -> str:
    """Remove maneuver_id and trailing underscores/dashes (case-insensitive)."""
    pattern = '^' + re.escape(maneuver_id) + r'[_\-]*'
    return re.sub(pattern, '', st_key, flags=re.IGNORECASE)


def clean_vehicle_names(vehicle_names: list[str]) -> dict[str, str]:
    """Remove common prefix/suffix and truncate to 34 chars."""
    cleaned = {}
    if len(vehicle_names) > 1:
        prefix = longest_common_prefix(vehicle_names)
        suffix = longest_common_suffix(vehicle_names)
    else:
        prefix = suffix = ""

    for name in vehicle_names:
        new_name = name
        if prefix:
            new_name = new_name[len(prefix):]
        if suffix and new_name.endswith(suffix):
            new_name = new_name[:-len(suffix)]

        new_name = new_name.strip("_- ")

        if len(new_name) > 33:
            new_name = new_name[:31] + "..."

        cleaned[name] = new_name
    return cleaned

# ============================================================================
# FIGURE CONVERSION - HIGH QUALITY
# ============================================================================

def figure_to_bytes(fig, is_plotly: bool = False) -> bytes:
    """
    Convert figure to bytes (HIGH QUALITY).
    Matplotlib: dpi=300
    Plotly: scale=2
    
    Returns bytes that can be processed by cv2.
    """
    try:
        if is_plotly:
            try:
                # Plotly HIGH QUALITY
                image_bytes = fig.to_image(format='png', width=854, height=586, scale=PLOTLY_SCALE)
                return image_bytes
            except Exception:
                # Fallback
                buf = io.BytesIO()
                fig.write_image(buf, format='png', width=854, height=586, scale=PLOTLY_SCALE)
                buf.seek(0)
                return buf.getvalue()
        else:
            # Matplotlib HIGH QUALITY
            buf = io.BytesIO()
            fig.savefig(buf, dpi=MATPLOTLIB_DPI, format='png', bbox_inches='tight', 
                       facecolor='white', edgecolor='none')
            buf.seek(0)
            plt.close(fig)
            return buf.getvalue()
    
    except Exception as e:
        st.warning(f"Figure conversion failed: {str(e)[:50]}")
        buf = io.BytesIO()
        fig_err, ax = plt.subplots(figsize=(12, 8), dpi=100)
        ax.text(0.5, 0.5, f"Error: {str(e)[:30]}", ha='center', va='center')
        ax.axis('off')
        fig_err.savefig(buf, dpi=100, format='png', bbox_inches='tight')
        plt.close(fig_err)
        buf.seek(0)
        return buf.getvalue()

# ============================================================================
# CV2 LEGEND DETECTION PIPELINE
# ============================================================================

def bytes_to_cv2(image_bytes: bytes) -> np.ndarray:
    """Convert image bytes to OpenCV format (BGR numpy array)."""
    nparr = cv2.imdecode(np.frombuffer(image_bytes, np.uint8), cv2.IMREAD_COLOR)
    if nparr is None:
        raise ValueError("Could not decode image from bytes")
    return nparr


def cv2_detect_legend_bbox(cv_img: np.ndarray, fallback_ratio: float = LEGEND_FALLBACK_RATIO) -> tuple:
    """
    Detect legend bounding box using OpenCV contour detection.
    
    Algorithm:
    1. Convert to grayscale
    2. Apply inverse binary threshold
    3. Find contours
    4. Look for legend-like contours (right side, reasonable size)
    5. Return bounding box or fallback to split ratio
    
    Args:
        cv_img: OpenCV image (BGR numpy array)
        fallback_ratio: Where to split if detection fails
    
    Returns:
        (x, y, w, h), height, width
    """
    h, w = cv_img.shape[:2]

    # Convert to grayscale
    gray = cv2.cvtColor(cv_img, cv2.COLOR_BGR2GRAY)
    
    # Inverse binary threshold - finds non-white areas (legend items)
    _, thresh = cv2.threshold(gray, 200, 255, cv2.THRESH_BINARY_INV)
    
    # Find contours
    contours, _ = cv2.findContours(thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

    # Look for legend contour
    legend_bbox = None
    for cnt in contours:
        x, y, cw, ch = cv2.boundingRect(cnt)
        # Legend is typically on right side, reasonable size
        if x > w * 0.6 and cw > 50 and ch > 50:
            legend_bbox = (x, y, cw, ch)
            break

    # Fallback: vertical split
    if legend_bbox is None:
        split_x = int(w * fallback_ratio)
        legend_bbox = (split_x, 0, w - split_x, h)

    return legend_bbox, h, w


def crop_and_separate(cv_img: np.ndarray, legend_bbox: tuple, 
                      left_padding: int = 30, right_padding: int = 30,
                      top_padding: int = 20, bottom_padding: int = 20,is_plotly=False) -> tuple:
    """
    Crop image into plot and legend with PADDING to show borders.
    
    LEGEND CROPPING: Increased padding to prevent legend cutoff
    Adjust these values if legend is still being cropped:
    - left_padding: Add space to left of legend (shows left border)
    - right_padding: Add space to right of legend (shows right border)  
    - top_padding: Add space above legend (shows top border)
    - bottom_padding: Add space below legend (shows bottom border)
    
    Args:
        cv_img: OpenCV image (BGR)
        legend_bbox: (x, y, w, h) from cv2.boundingRect
        left_padding: Pixels to add on left (show left border)
        right_padding: Pixels to add on right (show right border)
        top_padding: Pixels to add on top (show top border)
        bottom_padding: Pixels to add on bottom (show bottom border)
    
    Returns:
        (plot_bgr, legend_bgr) as OpenCV images
    """
    if is_plotly:
        left_padding = 0
        right_padding = 0
        top_padding = 0
        bottom_padding = 0
    h, w = cv_img.shape[:2]
    x, y, cw, ch = legend_bbox
    
    # Add padding to show legend borders AND prevent cutoff
    x_start = max(0, x - left_padding)
    x_end = min(w, x + cw + right_padding)
    y_start = max(0, y - top_padding)
    y_end = min(h, y + ch + bottom_padding)
    
    # Plot: everything left of legend (with padding)
    plot_bgr = cv_img[:, :x_start]
    
    # Legend: bounding box WITH padding to show borders
    legend_bgr = cv_img[y_start:y_end, x_start:x_end]
    
    return plot_bgr, legend_bgr


def cv2_to_pil(cv_img: np.ndarray) -> Image.Image:
    """Convert OpenCV BGR image to PIL Image (RGB)."""
    rgb = cv2.cvtColor(cv_img, cv2.COLOR_BGR2RGB)
    return Image.fromarray(rgb)


def pil_to_cv2(pil_img: Image.Image) -> np.ndarray:
    """Convert PIL Image to OpenCV BGR image."""
    rgb_arr = np.array(pil_img)
    bgr = cv2.cvtColor(rgb_arr, cv2.COLOR_RGB2BGR)
    return bgr

# ============================================================================
# HIGH QUALITY RESIZE WITH ASPECT RATIO
# ============================================================================

from PIL import Image

PLACEHOLDER_SIZE = (854, 586)

def respect(image: Image.Image,
            target_size: tuple = PLACEHOLDER_SIZE,
            bg_color: tuple = (255, 255, 255)) -> Image.Image:
    """
    Fit any image inside a fixed placeholder while maintaining aspect ratio.
    Always scales down if larger, scales up if smaller, then pads to target size.

    Args:
        image: PIL Image
        target_size: (width, height) of placeholder
        bg_color: RGB tuple for padding background

    Returns:
        PIL Image of exact target_size with aspect ratio preserved
    """
    target_w, target_h = target_size
    w, h = image.size

    # Compute scale factor to fit within target box
    scale = min(target_w / w, target_h / h)
    new_w, new_h = int(w * scale), int(h * scale)

    # Resize with high quality
    resized = image.resize((new_w, new_h), Image.LANCZOS)

    # Create canvas and paste centered
    canvas = Image.new("RGB", target_size, bg_color)
    offset = ((target_w - new_w) // 2, (target_h - new_h) // 2)
    canvas.paste(resized, offset)

    return canvas



def resend(image: Image.Image, target_width_px: int = LEGEND_SIZE[0], 
           target_height_px: int = LEGEND_SIZE[1], 
           bg_color: tuple = (255, 255, 255)) -> Image.Image:
    """
    Resize legend image to fit target box with aspect ratio maintained.
    
    Args:
        image: PIL Image
        target_width_px: Target width
        target_height_px: Target height
        bg_color: RGB tuple for padding
    
    Returns:
        PIL Image at target size with aspect ratio preserved
    """
    w, h = image.size
    scale = min(target_width_px / w, target_height_px / h)
    new_w = int(w * scale)
    new_h = int(h * scale)

    resized = image.resize((new_w, new_h), Image.LANCZOS)
    canvas = Image.new("RGB", (target_width_px, target_height_px), bg_color)
    offset_x = (target_width_px - new_w) // 2
    offset_y = (target_height_px - new_h) // 2
    canvas.paste(resized, (offset_x, offset_y))

    return canvas

# ============================================================================
# PRESENTATION BUILDING
# ============================================================================

def create_cover_slide(prs: Presentation, filename: str) -> None:
    """Create cover slide."""
    first_layout = prs.slide_layouts[0]
    first_slide = prs.slides.add_slide(first_layout)
    first_slide.placeholders[0].text = filename
    today_str = datetime.today().strftime("%d %B %Y")
    first_slide.placeholders[1].text = today_str


def create_overview_slide(prs: Presentation, selected: list, filename: str) -> None:
    """Create overview slide."""
    overview_layout = prs.slide_layouts[2]
    overview_slide = prs.slides.add_slide(overview_layout)
    overview_slide.placeholders[11].text = filename
    
    for i, maneuver_id in enumerate(selected[:8], start=1):
        title_idx = 18 + (i - 1) * 2
        num_idx = title_idx + 1
        overview_slide.placeholders[title_idx].text = maneuver_id
        overview_slide.placeholders[num_idx].text = f"{i:02d}"
    
    for j in range(len(selected) + 1, 9):
        title_idx = 18 + (j - 1) * 2
        num_idx = title_idx + 1
        for idx in (title_idx, num_idx):
            try:
                ph = overview_slide.placeholders[idx]
                ph.element.getparent().remove(ph.element)
            except IndexError:
                break


def update_vehicle_aliases(cleaned_map: dict, alias_map: list) -> dict:
    """Update vehicle aliases."""
    updated = {}
    for idx, (name, val) in enumerate(cleaned_map.items()):
        alias = alias_map[idx] if idx < len(alias_map) else name
        updated[name] = alias if alias != name else val
    return updated


def _format_table_cell_value(val) -> str:
    """Format value for table cell."""
    try:
        num = float(val)
        return f"{num:.4f}" if abs(num) < 0.01 else f"{num:.2f}"
    except (ValueError, TypeError):
        return str(val)


def _style_table_header(table, cols: int) -> None:
    """Apply header styling."""
    for c in range(cols):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_COLOR
        para = cell.text_frame.paragraphs[0]
        run = para.runs[0] if para.runs else para.add_run()
        run.font.size = TABLE_FONT_SIZE
        run.font.bold = True
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_cell_margins(cell)


def _style_table_data_cell(cell, value: str, is_numeric: bool = False) -> None:
    """Apply data cell styling."""
    cell.text = value
    para = cell.text_frame.paragraphs[0]
    if not para.runs:
        para.add_run()
    para.runs[0].font.size = TABLE_FONT_SIZE
    cell.text_frame.word_wrap = False
    _set_cell_margins(cell)
    if is_numeric:
        para.alignment = PP_ALIGN.CENTER


def _set_cell_margins(cell, margin_pt: int = 2) -> None:
    """Set cell margins."""
    cell.margin_left = Pt(margin_pt)
    cell.margin_top = Pt(0)
    cell.margin_bottom = Pt(0)
    cell.margin_right = Pt(margin_pt)


def insert_stowaway_tables(slide, maneuver_id: str, page_idx: int, 
                          table_mapping: dict) -> None:
    """Insert data tables into slide."""
    page_key = str(page_idx)
    if maneuver_id not in table_mapping or page_key not in table_mapping[maneuver_id]:
        return

    stowaway_dict = st.session_state.get("stowaways", {}).get(maneuver_id, {})
    imps = st.session_state.get("impostors", {}).get(maneuver_id, {}).get(f"Page_{page_idx}", [])
    
    for st_key in table_mapping[maneuver_id][page_key]:
        if st_key not in stowaway_dict:
            continue

        vehicle_data = stowaway_dict[st_key]
        rows, cols = len(vehicle_data) + 1, 2

        ph = slide.placeholders[18]
        graphic_frame = ph.insert_table(rows, cols)
        table = graphic_frame.table

        no_style_no_grid = '{21E7435A-961F-4814-916E-73D13E9FBCD9}'
        tbl = graphic_frame._element.graphic.graphicData.tbl
        if len(tbl) > 0:
            tbl[0][-1].text = no_style_no_grid

        table.cell(0, 0).text = "Vehicle"
        table.cell(0, 1).text = shorten_st_key(st_key, maneuver_id)
        _style_table_header(table, cols)

        cleaned_map = clean_vehicle_names(list(vehicle_data.keys()))
        cleaned = update_vehicle_aliases(cleaned_map, imps)
        
        for r, (veh, val) in enumerate(vehicle_data.items(), start=1):
            _style_table_data_cell(table.cell(r, 0), cleaned[veh])
            _style_table_data_cell(table.cell(r, 1), _format_table_cell_value(val), 
                                  is_numeric=True)

        total_width = graphic_frame.width
        table.columns[0].width = int(total_width * 0.7)
        table.columns[1].width = int(total_width * 0.3)

        total_height = graphic_frame.height
        natural_height = total_height / rows
        row_height = min(natural_height, DEFAULT_MAX_LEGEND_HEIGHT)
        for row in table.rows:
            row.height = row_height


def create_maneuver_slides(prs: Presentation, maneuver_id: str, figures: list,
                          cutfactor_1d: float, cutfactor_2d: float,
                          table_mapping: dict = None, slide_number: int = None,
                          progress_callback=None) -> None:
    """
    Create maneuver slides with cv2 legend detection pipeline.
    
    Pipeline:
    1. figure -> bytes (HIGH QUALITY: dpi=300, scale=2)
    2. bytes -> cv2 (BGR numpy array)
    3. cv2_detect_legend_bbox (cv2 contour detection)
    4. crop_and_separate (split into plot + legend WITH PADDING)
    5. cv2_to_pil (convert to PIL for resizing)
    6. respect/resend (PIL LANCZOS resize with aspect ratio)
    7. PIL -> bytes (back to bytes for PPTX)
    """
    if table_mapping is None:
        table_mapping = {}

    title_layout = prs.slide_layouts[3]
    title_slide = prs.slides.add_slide(title_layout)
    if slide_number is not None:
        title_slide.placeholders[14].text = f"{slide_number:02d}"
    title_slide.placeholders[15].text = maneuver_id

    for idx, figdata in enumerate(figures):
        if progress_callback:
            progress_callback(f"🎬 {maneuver_id} Page {idx+1}/{len(figures)}")
        
        try:
            # STEP 1: Render figure
            fig, split_ratio = figdata.render_page(
                page_data=figdata.page_data,
                vehicle_labels=st.session_state["impostors"].get(maneuver_id, {}).get(f"Page_{idx}", []),
                cutfactor_1d=cutfactor_1d,
                cutfactor_2d=cutfactor_2d
            )

            # STEP 2: Figure to HIGH QUALITY bytes
            is_plotly = not isinstance(fig, plt.Figure)
            image_bytes = figure_to_bytes(fig, is_plotly=is_plotly)
            
            # STEP 3: Bytes to OpenCV
            cv_img = bytes_to_cv2(image_bytes)
            
            # STEP 4: cv2 LEGEND DETECTION
            legend_bbox, h, w = cv2_detect_legend_bbox(cv_img, fallback_ratio=split_ratio)
            
            # STEP 5: CROP AND SEPARATE (with padding to show borders)
            plot_bgr, legend_bgr = crop_and_separate(
                cv_img, 
                legend_bbox,
                left_padding=LEGEND_PADDING['left'],
                right_padding=LEGEND_PADDING['right'],
                top_padding=LEGEND_PADDING['top'],
                bottom_padding=LEGEND_PADDING['bottom'],is_plotly
            )
            
            # STEP 6: CONVERT cv2 -> PIL
            plot_pil = cv2_to_pil(plot_bgr)
            legend_pil = cv2_to_pil(legend_bgr)
            
            # STEP 7: HIGH QUALITY RESIZE (LANCZOS + aspect ratio)
            # Different sizes for 1D vs 2D to avoid squeezing
            if is_plotly:
                # 2D Plotly plots - larger size
                plot_resized = respect(plot_pil, PLOT_SIZE_2D)
            else:
                # 1D Matplotlib plots - original size
                plot_resized = respect(plot_pil, PLOT_SIZE_1D)
            
            legend_resized = resend(legend_pil, *LEGEND_SIZE)
            
            # STEP 8: PIL -> BYTES (for PPTX)
            plot_bytes = io.BytesIO()
            plot_resized.save(plot_bytes, format='PNG')
            plot_bytes.seek(0)

            legend_bytes = io.BytesIO()
            legend_resized.save(legend_bytes, format='PNG')
            legend_bytes.seek(0)

            # Cleanup
            plot_pil.close()
            legend_pil.close()
            plot_resized.close()
            legend_resized.close()

            # STEP 9: ADD TO SLIDE
            slide_layout = prs.slide_layouts[4]
            slide = prs.slides.add_slide(slide_layout)
            
            try:
                slide.placeholders[11].text = f"{maneuver_id} - {figdata.page_data.get('title', f'Page{idx}')}"
            except KeyError:
                pass
            
            try:
                slide.placeholders[16].insert_picture(plot_bytes)
            except KeyError:
                slide.shapes.add_picture(plot_bytes, Inches(0.5), Inches(1.5),
                                        width=Inches(8.9), height=Inches(6.1))
            try:
                slide.placeholders[17].insert_picture(legend_bytes)
            except KeyError:
                slide.shapes.add_picture(legend_bytes, Inches(6.0), Inches(1.5))

            if maneuver_id in table_mapping and str(idx) in table_mapping[maneuver_id]:
                insert_stowaway_tables(slide, maneuver_id, idx, table_mapping)
            else:
                try:
                    ph = slide.placeholders[18]
                    ph.element.getparent().remove(ph.element)
                except IndexError:
                    pass

        except Exception as e:
            st.error(f"❌ {maneuver_id}:{idx} - {str(e)[:100]}")
            import traceback
            st.error(traceback.format_exc())
            continue

# ============================================================================
# STREAMLIT UI
# ============================================================================

def maneuver_selection(prefix: str = "") -> dict[str, bool]:
    """Get maneuver selection."""
    chosen = {}
    for m in st.session_state["tabs"].keys():
        chosen[m] = st.checkbox(f"Maneuver {m}", value=True, 
                               key=f"{prefix}chk_{m}")
    return chosen


def file_settings(prefix: str = "") -> tuple[str, str]:
    """Get file settings."""
    st.info("📁 **Streamlit Cloud:** Downloads automatically after export.")
    filename = st.text_input("Filename:", key=f"{prefix}ppt_filename", 
                            value="Export", max_chars=50)
    template_path = get_template_path()
    return filename, template_path


def export_settings(prefix: str = "") -> tuple[float, float]:
    """Get export settings."""
    with st.expander("⚙️ Settings", expanded=False):
        cutfactor_1d = st.number_input("Cut 1D", min_value=0.1, max_value=5.0,
                                      value=1.4, step=0.1, 
                                      key=f"{prefix}cutfactor_1d")
        cutfactor_2d = st.number_input("Cut 2D", min_value=0.1, max_value=10.0,
                                      value=7.5, step=0.1, 
                                      key=f"{prefix}cutfactor_2d")
    return cutfactor_1d, cutfactor_2d


@st.fragment
def pptX_tree(key=None):
    """Tree-based table selection UI."""
    st.subheader("📊 Add Tables")

    if "show_tree" not in st.session_state:
        st.session_state["show_tree"] = False

    if st.button("➕ Add Table", key="add_table_btn"):
        st.session_state["show_tree"] = not st.session_state["show_tree"]

    if not st.session_state["show_tree"]:
        return

    tree_data = [
        {"title": "cover_slide", "value": "cover"},
        {"title": "overview_slide", "value": "overview"},
    ]

    for m in st.session_state["tabs"].keys():
        man_node = {"title": f"Maneuver {m}", "value": f"man{m}", "children": []}
        for page_idx, figdata in enumerate(st.session_state["tabs"][m]["figures"], start=0):
            page_name = figdata.page_data.get("title", f"Page{page_idx}")
            page_node = {"title": page_name, "value": f"man{m}_page{page_idx}", 
                        "children": []}

            stowaway_dict = st.session_state.get("stowaways", {}).get(m)
            if stowaway_dict is None:
                stowaway_dict = st.session_state["tabs"][m].get("stowaway", {})

            if isinstance(stowaway_dict, dict):
                for option_key in stowaway_dict.keys():
                    page_node["children"].append({
                        "title": option_key,
                        "value": f"man{m}_page{page_idx}_{option_key}"
                    })
            elif isinstance(stowaway_dict, list):
                for param in stowaway_dict:
                    page_node["children"].append({
                        "title": param,
                        "value": f"man{m}_page{page_idx}_{param}"
                    })

            man_node["children"].append(page_node)
        tree_data.append(man_node)
    
    selected_values = st_ant_tree(
        treeData=tree_data,
        showSearch=True,
        treeCheckable=True,
        defaultValue=st.session_state.get("default_selected", []),
        key=key,
        max_height=600,
        placeholder="Search and select"
    )

    table_mapping = {}
    if selected_values:
        for val in selected_values:
            if "_page" in val:
                man_part, rest = val.split("_page", 1)
                page_part, stowaway_key = rest.split("_", 1)
                maneuver_id = man_part.replace("man", "", 1)
                page_no = page_part

                if maneuver_id not in table_mapping:
                    table_mapping[maneuver_id] = {}
                if page_no not in table_mapping[maneuver_id]:
                    table_mapping[maneuver_id][page_no] = []
                table_mapping[maneuver_id][page_no].append(stowaway_key)

    st.session_state["table_mapping"] = table_mapping


@st.fragment
def pptX_export():
    """Export to PowerPoint UI."""
    st.subheader("📥 Export Presentation")
    chosen = maneuver_selection(prefix="cloud_")
    filename, template_path = file_settings(prefix="cloud_")
    cutfactor_1d, cutfactor_2d = export_settings(prefix="cloud_")

    col1, col2 = st.columns([1, 1])
    
    with col1:
        export_btn = st.button("🚀 Generate", key="ppt_export", width="stretch")
    
    with col2:
        if st.session_state.get("pptx_ready", False):
            pptx_path = st.session_state.get("pptx_path", "")
            if os.path.exists(pptx_path):
                with open(pptx_path, "rb") as f:
                    st.download_button(
                        label="⬇️ Download",
                        data=f,
                        file_name=f"{filename}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        width="stretch"
                    )

    if export_btn:
        selected = [m for m, v in chosen.items() if v]
        
        if not selected:
            st.warning("Select at least one maneuver")
            return
        
        if not filename.strip():
            st.warning("Enter filename")
            return

        try:
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            prs = Presentation(template_path)
            temp_dir = get_temp_dir()

            status_text.text("📄 Cover...")
            create_cover_slide(prs, filename)
            progress_bar.progress(10)

            status_text.text("📊 Overview...")
            create_overview_slide(prs, selected, filename)
            progress_bar.progress(20)

            def update_progress(msg):
                status_text.text(msg)

            for i, maneuver_id in enumerate(selected, start=1):
                figures = st.session_state["tabs"][maneuver_id]["figures"]
                table_mapping = st.session_state.get("table_mapping", {})
                
                create_maneuver_slides(
                    prs, maneuver_id, figures,
                    cutfactor_1d, cutfactor_2d,
                    table_mapping=table_mapping,
                    slide_number=i,
                    progress_callback=update_progress
                )
                progress_bar.progress(20 + int((i / len(selected)) * 70))

            prs.slides.add_slide(prs.slide_layouts[9])
            ppt_path = os.path.join(temp_dir, f"{filename}.pptx")
            
            status_text.text("💾 Saving...")
            prs.save(ppt_path)
            progress_bar.progress(95)
            
            st.session_state["pptx_path"] = ppt_path
            st.session_state["pptx_ready"] = True
            
            progress_bar.progress(100)
            st.success(f"✅ Done! ({os.path.getsize(ppt_path) / (1024*1024):.1f}MB)")
            st.balloons()
            st.rerun()
        
        except Exception as e:
            st.error(f"❌ {str(e)[:150]}")


def pptX_tab():
    """Main interface."""
    tab1, tab2 = st.tabs(["📋 Tables", "📤 Export"])
    with tab1:
        pptX_tree(key="pptx_tree")
    with tab2:
        pptX_export()
