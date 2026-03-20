# -*- coding: utf-8 -*-
"""
Created on Tue Mar 17 10:13:51 2026

@tag: Xx_ScriptSniper_xX
@author: Albin

Optimized for Streamlit Cloud deployment
"""
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import cv2
import matplotlib.pyplot as plt
import streamlit as st
import shutil
import os
import re
import tempfile
import io
from datetime import datetime
from pathlib import Path
from st_ant_tree import st_ant_tree

# ============================================================================
# CONFIGURATION FOR STREAMLIT CLOUD
# ============================================================================

@st.cache_resource
def get_temp_dir():
    """Get temp directory for Streamlit Cloud."""
    temp_base = tempfile.gettempdir()
    temp_dir = os.path.join(temp_base, "streamlit_pptx_export")
    os.makedirs(temp_dir, exist_ok=True)
    return temp_dir


@st.cache_resource
def get_template_path() -> str:
    """Get template path from package directory."""
    base_dir = os.path.dirname(os.path.abspath(__file__))
    template = os.path.join(base_dir, "template.pptx")
    if not os.path.exists(template):
        raise FileNotFoundError(f"Template not found: {template}")
    return template


# Constants
DEFAULT_PLOT_SIZE = (854, 586)
DEFAULT_LEGEND_SIZE = (348, 344)
DEFAULT_MAX_LEGEND_HEIGHT = Cm(0.45)
TABLE_HEADER_COLOR = RGBColor(220, 220, 220)
TABLE_FONT_SIZE = Pt(10)
LEGEND_FALLBACK_RATIO = 0.75
DEFAULT_DPI = 150  # REDUCED from 300 for speed
TEMP_IMAGE_FOLDER = "ppt_temp_images"
CHUNK_SIZE = 5242880  # 5MB for file downloads

# ============================================================================
# UTILITY FUNCTIONS
# ============================================================================

def longest_common_prefix(strings: list[str]) -> str:
    """Find longest common prefix in a list of strings."""
    if not strings:
        return ""
    s1, s2 = min(strings), max(strings)
    for i, c in enumerate(s1):
        if c != s2[i]:
            return s1[:i]
    return s1


def longest_common_suffix(strings: list[str]) -> str:
    """Find longest common suffix in a list of strings."""
    if not strings:
        return ""
    rev = [s[::-1] for s in strings]
    return longest_common_prefix(rev)[::-1]


def shorten_st_key(st_key: str, maneuver_id: str) -> str:
    """Remove maneuver_id and trailing underscores/dashes (case-insensitive)."""
    pattern = '^' + re.escape(maneuver_id) + r'[_\-]*'
    return re.sub(pattern, '', st_key, flags=re.IGNORECASE)


def clean_vehicle_names(vehicle_names: list[str]) -> dict[str, str]:
    """Remove common prefix/suffix and truncate to 34 chars."""
    cleaned = {}
    if len(vehicle_names) > 1:
        prefix = longest_common_prefix(vehicle_names)
        suffix = longest_common_suffix(vehicle_names)
    else:
        prefix = suffix = ""

    for name in vehicle_names:
        new_name = name
        if prefix:
            new_name = new_name[len(prefix):]
        if suffix and new_name.endswith(suffix):
            new_name = new_name[:-len(suffix)]

        new_name = new_name.strip("_- ")

        if len(new_name) > 33:
            new_name = new_name[:31] + "..."

        cleaned[name] = new_name
    return cleaned

# ============================================================================
# IMAGE PROCESSING FUNCTIONS (WITH OPTIMIZATION)
# ============================================================================

def optimize_image(image_path: str, quality: int = 85) -> None:
    """Optimize image file size after saving using PIL."""
    try:
        img = Image.open(image_path)
        # For PNG, use optimize flag; for JPEG, use quality
        if image_path.lower().endswith('.png'):
            img.save(image_path, 'PNG', optimize=True)
        else:
            img.save(image_path, quality=quality)
        img.close()
    except Exception as e:
        pass  # Silent fail - optimization not critical


def respect(image: Image.Image, target_size: tuple = DEFAULT_PLOT_SIZE, 
            bg_color: tuple = (255, 255, 255)) -> Image.Image:
    """Resize image proportionally and pad to target_size."""
    target_w, target_h = target_size
    w, h = image.size
    scale = min(target_w / w, target_h / h)
    new_w, new_h = int(w * scale), int(h * scale)

    resized = image.resize((new_w, new_h), Image.LANCZOS)
    canvas = Image.new("RGB", target_size, bg_color)
    offset = ((target_w - new_w) // 2, (target_h - new_h) // 2)
    canvas.paste(resized, offset)
    return canvas


def resend(image: Image.Image, target_width_px: int = DEFAULT_LEGEND_SIZE[0],
           target_height_px: int = DEFAULT_LEGEND_SIZE[1], 
           bg_color: tuple = (255, 255, 255)) -> Image.Image:
    """Resize legend image to fit target box while maintaining aspect ratio."""
    w, h = image.size
    scale = min(target_width_px / w, target_height_px / h)
    new_w, new_h = int(w * scale), int(h * scale)

    resized = image.resize((new_w, new_h), Image.LANCZOS)
    canvas = Image.new("RGB", (target_width_px, target_height_px), bg_color)
    offset_x = (target_width_px - new_w) // 2
    offset_y = (target_height_px - new_h) // 2
    canvas.paste(resized, (offset_x, offset_y))
    return canvas


def detect_legend_bbox(base_img_path: str, fallback_ratio: float = LEGEND_FALLBACK_RATIO
                      ) -> tuple:
    """Detect legend bounding box using grayscale + threshold."""
    cv_img = cv2.imread(base_img_path)
    if cv_img is None:
        raise FileNotFoundError(f"Could not read image: {base_img_path}")
    
    h, w = cv_img.shape[:2]
    gray = cv2.cvtColor(cv_img, cv2.COLOR_BGR2GRAY)
    _, thresh = cv2.threshold(gray, 200, 255, cv2.THRESH_BINARY_INV)
    contours, _ = cv2.findContours(thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

    legend_bbox = None
    for cnt in contours:
        x, y, cw, ch = cv2.boundingRect(cnt)
        if x > w * 0.6 and cw > 50 and ch > 50:
            legend_bbox = (x, y, cw, ch)
            break

    if legend_bbox is None:
        split_x = int(w * fallback_ratio)
        legend_bbox = (split_x, 0, w - split_x, h)

    return legend_bbox, h, w

# ============================================================================
# PRESENTATION BUILDING FUNCTIONS
# ============================================================================

def create_cover_slide(prs: Presentation, filename: str) -> None:
    """Create cover slide with filename and date."""
    first_layout = prs.slide_layouts[0]
    first_slide = prs.slides.add_slide(first_layout)
    first_slide.placeholders[0].text = filename
    today_str = datetime.today().strftime("%d %B %Y")
    first_slide.placeholders[1].text = today_str


def create_overview_slide(prs: Presentation, selected: list, filename: str) -> None:
    """Create overview slide with selected maneuvers."""
    overview_layout = prs.slide_layouts[2]
    overview_slide = prs.slides.add_slide(overview_layout)
    overview_slide.placeholders[11].text = filename
    
    for i, maneuver_id in enumerate(selected[:8], start=1):
        title_idx = 18 + (i - 1) * 2
        num_idx = title_idx + 1
        overview_slide.placeholders[title_idx].text = maneuver_id
        overview_slide.placeholders[num_idx].text = f"{i:02d}"
    
    for j in range(len(selected) + 1, 9):
        title_idx = 18 + (j - 1) * 2
        num_idx = title_idx + 1
        for idx in (title_idx, num_idx):
            try:
                ph = overview_slide.placeholders[idx]
                ph.element.getparent().remove(ph.element)
            except IndexError:
                break


def update_vehicle_aliases(cleaned_map: dict, alias_map: list) -> dict:
    """Update vehicle aliases from alias_map."""
    updated = {}
    for idx, (name, val) in enumerate(cleaned_map.items()):
        alias = alias_map[idx] if idx < len(alias_map) else name
        updated[name] = alias if alias != name else val
    return updated


def _format_table_cell_value(val) -> str:
    """Format value for table cell."""
    try:
        num = float(val)
        return f"{num:.4f}" if abs(num) < 0.01 else f"{num:.2f}"
    except (ValueError, TypeError):
        return str(val)


def _style_table_header(table, cols: int) -> None:
    """Apply header styling to table."""
    for c in range(cols):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_COLOR
        para = cell.text_frame.paragraphs[0]
        run = para.runs[0] if para.runs else para.add_run()
        run.font.size = TABLE_FONT_SIZE
        run.font.bold = True
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_cell_margins(cell)


def _style_table_data_cell(cell, value: str, is_numeric: bool = False) -> None:
    """Apply data cell styling."""
    cell.text = value
    para = cell.text_frame.paragraphs[0]
    if not para.runs:
        para.add_run()
    para.runs[0].font.size = TABLE_FONT_SIZE
    cell.text_frame.word_wrap = False
    _set_cell_margins(cell)
    if is_numeric:
        para.alignment = PP_ALIGN.CENTER


def _set_cell_margins(cell, margin_pt: int = 2) -> None:
    """Set consistent cell margins."""
    cell.margin_left = Pt(margin_pt)
    cell.margin_top = Pt(0)
    cell.margin_bottom = Pt(0)
    cell.margin_right = Pt(margin_pt)


def insert_stowaway_tables(slide, maneuver_id: str, page_idx: int, 
                          table_mapping: dict) -> None:
    """Insert data tables into slide."""
    page_key = str(page_idx)
    if maneuver_id not in table_mapping or page_key not in table_mapping[maneuver_id]:
        return

    stowaway_dict = st.session_state.get("stowaways", {}).get(maneuver_id, {})
    imps = st.session_state.get("impostors", {}).get(maneuver_id, {}).get(f"Page_{page_idx}", [])
    
    for st_key in table_mapping[maneuver_id][page_key]:
        if st_key not in stowaway_dict:
            continue

        vehicle_data = stowaway_dict[st_key]
        rows, cols = len(vehicle_data) + 1, 2

        ph = slide.placeholders[18]
        graphic_frame = ph.insert_table(rows, cols)
        table = graphic_frame.table

        # Apply "No Style, No Grid" table style
        no_style_no_grid = '{21E7435A-961F-4814-916E-73D13E9FBCD9}'
        tbl = graphic_frame._element.graphic.graphicData.tbl
        if len(tbl) > 0:
            tbl[0][-1].text = no_style_no_grid

        # Header row
        table.cell(0, 0).text = "Vehicle"
        table.cell(0, 1).text = shorten_st_key(st_key, maneuver_id)
        _style_table_header(table, cols)

        # Data rows
        cleaned_map = clean_vehicle_names(list(vehicle_data.keys()))
        cleaned = update_vehicle_aliases(cleaned_map, imps)
        
        for r, (veh, val) in enumerate(vehicle_data.items(), start=1):
            _style_table_data_cell(table.cell(r, 0), cleaned[veh])
            _style_table_data_cell(table.cell(r, 1), _format_table_cell_value(val), 
                                  is_numeric=True)

        # Column widths
        total_width = graphic_frame.width
        table.columns[0].width = int(total_width * 0.7)
        table.columns[1].width = int(total_width * 0.3)

        # Row heights
        total_height = graphic_frame.height
        natural_height = total_height / rows
        row_height = min(natural_height, DEFAULT_MAX_LEGEND_HEIGHT)
        for row in table.rows:
            row.height = row_height


def create_maneuver_slides(prs: Presentation, maneuver_id: str, figures: list,
                          cutfactor_1d: float, cutfactor_2d: float, temp_dir: str,
                          table_mapping: dict = None, slide_number: int = None,
                          progress_callback=None) -> None:
    """Create slides for a maneuver with figures and tables."""
    if table_mapping is None:
        table_mapping = {}

    title_layout = prs.slide_layouts[3]
    title_slide = prs.slides.add_slide(title_layout)
    if slide_number is not None:
        title_slide.placeholders[14].text = f"{slide_number:02d}"
    title_slide.placeholders[15].text = maneuver_id

    for idx, figdata in enumerate(figures):
        if progress_callback:
            progress_callback(f"Processing {maneuver_id} - Page {idx+1}/{len(figures)}")
        
        try:
            fig, split_ratio = figdata.render_page(
                page_data=figdata.page_data,
                vehicle_labels=st.session_state["impostors"].get(maneuver_id, {}).get(f"Page_{idx}", []),
                cutfactor_1d=cutfactor_1d,
                cutfactor_2d=cutfactor_2d
            )

            base_img_path = os.path.join(temp_dir, f"{maneuver_id}_Page{idx}_full.png")
            
            # Save figure based on type - FIXED
            if isinstance(fig, plt.Figure):
                # Matplotlib figure - save without optimize parameter
                fig.savefig(base_img_path, dpi=DEFAULT_DPI, bbox_inches="tight", format='png')
                plt.close(fig)
                # Optimize AFTER saving
                optimize_image(base_img_path)
            else:
                # Plotly or other figure type
                fig.write_image(base_img_path, scale=1)
                # Optimize after save
                optimize_image(base_img_path)
                
                img = Image.open(base_img_path)
                w, h = img.size
                margin = 50
                figcrop = img.crop((margin, 0, w, h))
                figcrop.save(base_img_path)
                img.close()

            # Legend detection and splitting
            legend_bbox, h, w = detect_legend_bbox(base_img_path, fallback_ratio=split_ratio)
            x, y, cw, ch = legend_bbox

            img = Image.open(base_img_path)
            plot_crop = img.crop((0, 0, x, h))
            legend_crop = img.crop((x, y, x + cw, y + ch))

            # Resize images
            plot_resized = respect(plot_crop, DEFAULT_PLOT_SIZE)
            plot_path = os.path.join(temp_dir, f"{maneuver_id}_Page{idx}_plot.png")
            plot_resized.save(plot_path)
            
            legend_resized = resend(legend_crop, *DEFAULT_LEGEND_SIZE)
            legend_path = os.path.join(temp_dir, f"{maneuver_id}_Page{idx}_legend.png")
            legend_resized.save(legend_path)

            # Optimize resized images
            optimize_image(plot_path)
            optimize_image(legend_path)

            img.close()
            plot_crop.close()
            legend_crop.close()
            plot_resized.close()
            legend_resized.close()

            # Create slide
            slide_layout = prs.slide_layouts[4]
            slide = prs.slides.add_slide(slide_layout)
            
            try:
                title = f"{maneuver_id} - {figdata.page_data.get('title', f'Page{idx}')}"
                slide.placeholders[11].text = title
            except KeyError:
                pass

            # Insert plot
            try:
                slide.placeholders[16].insert_picture(plot_path)
            except KeyError:
                slide.shapes.add_picture(plot_path, Inches(0.5), Inches(1.5),
                                        width=Inches(8.9), height=Inches(6.1))

            # Insert legend
            try:
                slide.placeholders[17].insert_picture(legend_path)
            except KeyError:
                slide.shapes.add_picture(legend_path, Inches(6.0), Inches(1.5))

            # Insert tables if applicable
            if maneuver_id in table_mapping and str(idx) in table_mapping[maneuver_id]:
                insert_stowaway_tables(slide, maneuver_id, idx, table_mapping)
            else:
                try:
                    ph = slide.placeholders[18]
                    ph.element.getparent().remove(ph.element)
                except IndexError:
                    pass

        except Exception as e:
            st.error(f"Error processing {maneuver_id} Page {idx}: {str(e)}")
            continue

# ============================================================================
# STREAMLIT UI FUNCTIONS
# ============================================================================

def maneuver_selection(prefix: str = "") -> dict[str, bool]:
    """Get maneuver selection from checkboxes."""
    chosen = {}
    for m in st.session_state["tabs"].keys():
        chosen[m] = st.checkbox(f"Maneuver {m}", value=True, 
                               key=f"{prefix}chk_{m}")
    return chosen


def file_settings(prefix: str = "") -> tuple[str, str, str]:
    """Get file settings - cloud-optimized."""
    st.info("📁 **For Streamlit Cloud:** Files are saved temporarily. Download immediately.")
    
    filename = st.text_input("Filename (no extension):", key=f"{prefix}ppt_filename", 
                            value="Export", max_chars=50)
    template_path = get_template_path()
    temp_dir = get_temp_dir()
    
    return temp_dir, filename, template_path


def export_settings(prefix: str = "") -> tuple[float, float]:
    """Get export settings from UI."""
    with st.expander("⚙️ Export Settings", expanded=False):
        cutfactor_1d = st.number_input("Cut factor 1D", min_value=0.1, max_value=5.0,
                                      value=1.4, step=0.1, 
                                      key=f"{prefix}cutfactor_1d")
        cutfactor_2d = st.number_input("Cut factor 2D", min_value=0.1, max_value=10.0,
                                      value=7.5, step=0.1, 
                                      key=f"{prefix}cutfactor_2d")
    return cutfactor_1d, cutfactor_2d


@st.fragment
def pptX_tree(key=None):
    """Tree-based table selection UI."""
    st.subheader("📊 Add Tables to Export")

    if "show_tree" not in st.session_state:
        st.session_state["show_tree"] = False

    if st.button("➕ Add Table", key="add_table_btn"):
        st.session_state["show_tree"] = not st.session_state["show_tree"]

    if not st.session_state["show_tree"]:
        return

    tree_data = [
        {"title": "cover_slide", "value": "cover"},
        {"title": "overview_slide", "value": "overview"},
    ]

    for m in st.session_state["tabs"].keys():
        man_node = {"title": f"Maneuver {m}", "value": f"man{m}", "children": []}
        for page_idx, figdata in enumerate(st.session_state["tabs"][m]["figures"], start=0):
            page_name = figdata.page_data.get("title", f"Page{page_idx}")
            page_node = {"title": page_name, "value": f"man{m}_page{page_idx}", 
                        "children": []}

            stowaway_dict = st.session_state.get("stowaways", {}).get(m)
            if stowaway_dict is None:
                stowaway_dict = st.session_state["tabs"][m].get("stowaway", {})

            if isinstance(stowaway_dict, dict):
                for option_key in stowaway_dict.keys():
                    page_node["children"].append({
                        "title": option_key,
                        "value": f"man{m}_page{page_idx}_{option_key}"
                    })
            elif isinstance(stowaway_dict, list):
                for param in stowaway_dict:
                    page_node["children"].append({
                        "title": param,
                        "value": f"man{m}_page{page_idx}_{param}"
                    })

            man_node["children"].append(page_node)
        tree_data.append(man_node)
    
    selected_values = st_ant_tree(
        treeData=tree_data,
        showSearch=True,
        treeCheckable=True,
        defaultValue=st.session_state.get("default_selected", []),
        key=key,
        max_height=600,
        placeholder="Search and select"
    )

    table_mapping = {}
    if selected_values:
        for val in selected_values:
            if "_page" in val:
                man_part, rest = val.split("_page", 1)
                page_part, stowaway_key = rest.split("_", 1)
                maneuver_id = man_part.replace("man", "", 1)
                page_no = page_part

                if maneuver_id not in table_mapping:
                    table_mapping[maneuver_id] = {}
                if page_no not in table_mapping[maneuver_id]:
                    table_mapping[maneuver_id][page_no] = []
                table_mapping[maneuver_id][page_no].append(stowaway_key)

    st.session_state["table_mapping"] = table_mapping


@st.fragment
def pptX_export():
    """Export to PowerPoint UI - Cloud optimized."""
    st.subheader("📥 Export PowerPoint Presentation")
    chosen = maneuver_selection(prefix="cloud_")
    directory, filename, template_path = file_settings(prefix="cloud_")
    cutfactor_1d, cutfactor_2d = export_settings(prefix="cloud_")

    col1, col2 = st.columns([1, 1])
    
    with col1:
        export_btn = st.button("🚀 Generate PPT", key="ppt_export", use_container_width=True)
    
    with col2:
        if st.session_state.get("pptx_ready", False):
            with open(st.session_state.get("pptx_path", ""), "rb") as f:
                st.download_button(
                    label="⬇️ Download PPT",
                    data=f,
                    file_name=f"{filename}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )

    if export_btn:
        selected = [m for m, v in chosen.items() if v]
        
        if not selected:
            st.warning("⚠️ Please select at least one maneuver.")
            return
        
        if not filename or not filename.strip():
            st.warning("⚠️ Please provide a filename.")
            return

        try:
            # Progress tracking
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            prs = Presentation(template_path)
            os.makedirs(directory, exist_ok=True)

            status_text.text("📄 Creating cover slide...")
            create_cover_slide(prs, filename)
            progress_bar.progress(10)

            status_text.text("📊 Creating overview slide...")
            create_overview_slide(prs, selected, filename)
            progress_bar.progress(20)

            def update_progress(msg):
                status_text.text(msg)

            for i, maneuver_id in enumerate(selected, start=1):
                figures = st.session_state["tabs"][maneuver_id]["figures"]
                table_mapping = st.session_state.get("table_mapping", {})
                
                create_maneuver_slides(
                    prs, maneuver_id, figures,
                    cutfactor_1d, cutfactor_2d, directory,
                    slide_number=i,
                    table_mapping=table_mapping,
                    progress_callback=update_progress
                )
                progress_bar.progress(20 + int((i / len(selected)) * 70))

            prs.slides.add_slide(prs.slide_layouts[9])
            ppt_path = os.path.join(directory, f"{filename}.pptx")
            
            status_text.text("💾 Saving presentation...")
            prs.save(ppt_path)
            progress_bar.progress(95)
            
            # Store path in session state for download
            st.session_state["pptx_path"] = ppt_path
            st.session_state["pptx_ready"] = True
            
            progress_bar.progress(100)
            st.success(f"✅ PPT generated successfully!")
            st.balloons()
            
            # Auto-rerun to show download button
            st.rerun()
        
        except Exception as e:
            st.error(f"❌ Export failed: {str(e)}")


def pptX_tab():
    """Main presentation export interface."""
    tab1, tab2 = st.tabs(["📋 Add Tables", "📤 Export"])
    with tab1:
        pptX_tree(key="pptx_tree")
    with tab2:
        pptX_export()
